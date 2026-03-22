from __future__ import annotations

import json
import os
import re
import subprocess
import sys
from pathlib import Path

from tools import convert_office_to_pdf, pdf_to_pngs, pdf_to_text
from workflows.common.tools import (
    ensure_dir,
    excerpt,
    get_stage_batch_size,
    normalize_outline_slide_sections,
    parse_bool,
    parse_outline_slide_sections,
    render_slide_batch_markdown,
    write_text,
)


def _load_json_like(raw: object, default):
    if raw is None:
        return default
    if isinstance(raw, (list, dict)):
        return raw
    text = str(raw).strip()
    if not text:
        return default
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        return default


def _sanitize_field_text(raw: object, *, tag_name: str) -> str:
    text = str(raw or "").strip()
    text = re.sub(rf"</?{re.escape(tag_name)}>", "", text, flags=re.IGNORECASE).strip()

    fenced_match = re.search(r"```(?:python|markdown)?\s*(.*?)```", text, flags=re.DOTALL | re.IGNORECASE)
    if fenced_match:
        text = fenced_match.group(1).strip()

    tagged_match = re.search(
        rf"<{re.escape(tag_name)}>\s*(.*?)\s*</{re.escape(tag_name)}>",
        text,
        flags=re.DOTALL | re.IGNORECASE,
    )
    if tagged_match:
        text = tagged_match.group(1).strip()

    return text.strip()


def _is_unset_placeholder(text: str) -> bool:
    normalized = text.strip().lower()
    return normalized in {"", "(not set yet)", "not set yet"}


def _extract_slide_sections(html_text: str) -> list[str]:
    pattern = re.compile(r"(?is)<section\b[^>]*class=[\"'][^\"']*\bslide\b[^\"']*[\"'][^>]*>.*?</section>")
    return [match.group(0).strip() for match in pattern.finditer(html_text or "")]


def _render_current_batch_html(html_text: str, start_index: int, end_index: int) -> str:
    sections = _extract_slide_sections(html_text)
    if not sections:
        return str(html_text or "")
    selected = sections[start_index:end_index]
    return "\n\n".join(selected).strip()


def _candidate_python_bins(attrs: dict[str, object]) -> list[str]:
    candidates: list[str] = []
    for value in [
        attrs.get("python_bin"),
        sys.executable,
        str(Path(__file__).resolve().parents[3] / ".venv" / "bin" / "python"),
    ]:
        python_bin = str(value or "").strip()
        if python_bin and python_bin not in candidates:
            candidates.append(python_bin)
    return candidates


def _python_supports_recreation(python_bin: str) -> bool:
    if not python_bin or not Path(python_bin).exists():
        return False
    completed = subprocess.run(
        [python_bin, "-c", "import PIL, pptx"],
        text=True,
        capture_output=True,
        check=False,
    )
    return completed.returncode == 0


def _select_recreation_python(attrs: dict[str, object]) -> str:
    candidates = _candidate_python_bins(attrs)
    for python_bin in candidates:
        if _python_supports_recreation(python_bin):
            return python_bin
    return candidates[0] if candidates else sys.executable


def run_recreation_script(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    script_text = _sanitize_field_text(attrs.get("python_pptx_script", ""), tag_name="python_pptx_script")
    notes_text = _sanitize_field_text(attrs.get("recreation_notes_markdown", ""), tag_name="recreation_notes_markdown")
    script_path = Path(str(attrs["recreated_script_path"]))
    notes_file_path = Path(str(attrs["recreation_notes_path"]))
    if _is_unset_placeholder(script_text) and script_path.exists():
        script_text = script_path.read_text(encoding="utf-8", errors="ignore")
    if _is_unset_placeholder(notes_text) and notes_file_path.exists():
        notes_text = notes_file_path.read_text(encoding="utf-8", errors="ignore")
    script_path = Path(write_text(script_path, script_text))
    notes_path = write_text(notes_file_path, notes_text)
    pptx_path = Path(str(attrs["recreated_pptx_path"]))
    expected_pdf_path = Path(str(attrs["recreated_validation_pdf_path"]))
    preview_dir = ensure_dir(attrs["recreated_preview_dir"])
    python_bin = _select_recreation_python(attrs)

    success = False
    run_stdout = ""
    run_stderr = ""
    recreated_pdf_text = ""
    preview_paths: list[Path] = []
    error_text = ""

    env = os.environ.copy()
    env["OHNOPPT_OUTPUT_PPTX"] = str(pptx_path)

    try:
        completed = subprocess.run(
            [python_bin, str(script_path)],
            text=True,
            capture_output=True,
            check=True,
            env=env,
            cwd=str(script_path.parent),
        )
        run_stdout = completed.stdout
        run_stderr = completed.stderr
        if pptx_path.exists():
            converted_pdf = convert_office_to_pdf(pptx_path, output_dir=pptx_path.parent)
            if converted_pdf != expected_pdf_path and converted_pdf.exists():
                expected_pdf_path.write_bytes(converted_pdf.read_bytes())
        if expected_pdf_path.exists():
            preview_paths = pdf_to_pngs(expected_pdf_path, preview_dir / "slide", dpi=144)
            text_path = str(expected_pdf_path).replace(".pdf", ".txt")
            pdf_to_text(expected_pdf_path, output_txt=text_path)
            recreated_pdf_text = Path(text_path).read_text(encoding="utf-8", errors="ignore")
        success = pptx_path.exists() and expected_pdf_path.exists()
    except Exception as exc:
        error_text = f"{type(exc).__name__}: {exc}"

    execution_report = "\n".join(
        [
            "# Recreation Execution Report",
            "",
            f"- Script: {script_path}",
            f"- Python: {python_bin}",
            f"- PPTX: {pptx_path}",
            f"- Validation PDF: {expected_pdf_path}",
            f"- Preview dir: {preview_dir}",
            f"- Success: {success}",
            f"- Error: {error_text or '(none)'}",
            "",
            "## stdout",
            excerpt(run_stdout, limit=5000) or "(empty)",
            "",
            "## stderr",
            excerpt(run_stderr, limit=5000) or "(empty)",
        ]
    )
    recreated_pptx_summary = "\n".join(
        [
            "# Recreated PPTX Summary",
            "",
            f"- PPTX path: {pptx_path}",
            f"- Validation PDF path: {expected_pdf_path}",
            f"- Preview dir: {preview_dir}",
            f"- Recreation success: {success}",
            "",
            "## Validation PDF Text Excerpt",
            excerpt(recreated_pdf_text, limit=5000) if recreated_pdf_text else "(not available)",
            "",
            "## Execution Report",
            execution_report,
        ]
    )
    return {
        "recreated_script_path": str(script_path),
        "recreation_notes_path": notes_path,
        "python_pptx_script": script_text,
        "recreation_notes_markdown": notes_text,
        "python_pptx_script_excerpt": excerpt(script_text, limit=14000),
        "pptx_recreation_success": success,
        "recreated_pptx_path": str(pptx_path),
        "recreated_validation_pdf_path": str(expected_pdf_path),
        "recreated_preview_dir": str(preview_dir),
        "recreated_pptx_summary_markdown": recreated_pptx_summary,
    }


def assemble_pptx_script(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    batch_outputs = _load_json_like(attrs.get("pptx_batch_outputs_json"), [])
    batch_outputs = sorted(batch_outputs, key=lambda item: int(item.get("batch_index") or 0))
    sections = _load_json_like(attrs.get("pptx_batch_sections_json"), [])

    slide_numbers = [int(item.get("number") or 0) for item in sections if int(item.get("number") or 0) > 0]
    fragment_blocks: list[str] = []
    notes_blocks: list[str] = []

    for item in batch_outputs:
        fragment = _sanitize_field_text(item.get("python_pptx_batch_fragment", ""), tag_name="python_pptx_batch_fragment")
        notes = _sanitize_field_text(item.get("pptx_batch_notes_markdown", ""), tag_name="pptx_batch_notes_markdown")
        batch_index = int(item.get("batch_index") or 0)
        if fragment:
            fragment_blocks.append(f"# ===== Batch {batch_index} =====\n{fragment.strip()}")
        if notes:
            notes_blocks.append(f"# Batch {batch_index}\n\n{notes.strip()}")

    script_text = "\n\n".join(
        part
        for part in [
            "#!/usr/bin/env python3",
            "from __future__ import annotations",
            "",
            "import os",
            "from pathlib import Path",
            "",
            "from PIL import Image",
            "from pptx import Presentation",
            "from pptx.dml.color import RGBColor",
            "from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE",
            "from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR",
            "from pptx.shapes.shapetree import SlideShapes",
            "from pptx.util import Inches, Pt, Cm",
            "",
            f'OUTPUT_PPTX = os.getenv("OHNOPPT_OUTPUT_PPTX", r"{attrs.get("recreated_pptx_path", "")}")',
            "",
            "def _compat_add_line(self, x1, y1, x2, y2):",
            "    return self.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)",
            "",
            "if not hasattr(SlideShapes, \"add_line\"):",
            "    SlideShapes.add_line = _compat_add_line",
            "",
            "\n\n".join(fragment_blocks).strip(),
            "",
            "def build_presentation():",
            "    prs = Presentation()",
            "    prs.slide_width = Inches(13.333333)",
            "    prs.slide_height = Inches(7.5)",
            *[f"    build_slide_{number}(prs)" for number in slide_numbers],
            "    return prs",
            "",
            "def main():",
            "    prs = build_presentation()",
            "    out = Path(OUTPUT_PPTX)",
            "    out.parent.mkdir(parents=True, exist_ok=True)",
            "    prs.save(str(out))",
            '    print(f"Saved PPTX to: {out}")',
            "",
            'if __name__ == "__main__":',
            "    main()",
            "",
        ]
        if str(part).strip()
    ).strip() + "\n"

    notes_text = "\n\n".join(notes_blocks).strip() or "# Recreation Notes\n\n(no batch notes)\n"
    script_path = Path(write_text(attrs["recreated_script_path"], script_text))
    notes_path = write_text(attrs["recreation_notes_path"], notes_text)
    return {
        "recreated_script_path": str(script_path),
        "recreation_notes_path": notes_path,
        "python_pptx_script": script_text,
        "recreation_notes_markdown": notes_text,
        "python_pptx_script_excerpt": excerpt(script_text, limit=14000),
    }


def pack_pptx_review_for_human(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    return {
        "stage": "pptx_recreation",
        "artifact_paths": (
            f"pptx: {attrs.get('recreated_pptx_path')}\n"
            f"validation_pdf: {attrs.get('recreated_validation_pdf_path')}\n"
            f"script: {attrs.get('recreated_script_path')}"
        ),
        "artifact_excerpt": excerpt(attrs.get("python_pptx_script"), limit=5000),
    }


def prepare_pptx_batch_plan(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    sections = normalize_outline_slide_sections(
        parse_outline_slide_sections(attrs.get("outline_markdown", "")),
        page_budget=attrs.get("page_budget"),
    )
    batch_size = get_stage_batch_size(default=10)
    batch_count = max(1, (len(sections) + batch_size - 1) // batch_size) if sections else 1
    return {
        "pptx_batch_sections_json": json.dumps(sections, ensure_ascii=False),
        "pptx_batch_size": batch_size,
        "pptx_batch_count": batch_count,
        "pptx_batch_index": 0,
        "pptx_batch_outputs_json": "[]",
        "pptx_batch_iteration_ready": False,
    }


def prepare_pptx_batch_iteration(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    sections = _load_json_like(attrs.get("pptx_batch_sections_json"), [])
    batch_size = int(attrs.get("pptx_batch_size") or get_stage_batch_size(default=10))
    batch_count = int(attrs.get("pptx_batch_count") or 1)
    batch_index = int(attrs.get("pptx_batch_index") or 0)
    html_content = str(attrs.get("html_content", ""))

    if sections:
        start = batch_index * batch_size
        end = min(start + batch_size, len(sections))
        current_sections = sections[start:end]
        current_batch_markdown = render_slide_batch_markdown(
            current_sections,
            batch_index=batch_index + 1,
            batch_count=batch_count,
        )
        current_batch_html_content = _render_current_batch_html(html_content, start, end)
        batch_label = (
            f"batch {batch_index + 1}/{batch_count} | slides "
            f"{current_sections[0]['number']}-{current_sections[-1]['number']}"
            if current_sections
            else f"batch {batch_index + 1}/{batch_count}"
        )
    else:
        current_batch_markdown = str(attrs.get("outline_markdown", ""))
        current_batch_html_content = html_content
        batch_label = "batch 1/1 | full deck"

    return {
        "pptx_current_batch_markdown": current_batch_markdown,
        "pptx_current_batch_html_content": current_batch_html_content,
        "pptx_batch_label": batch_label,
    }


def finalize_pptx_batch_iteration(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    batch_index = int(attrs.get("pptx_batch_index") or 0)
    batch_count = int(attrs.get("pptx_batch_count") or 1)
    next_index = batch_index + 1
    outputs = _load_json_like(attrs.get("pptx_batch_outputs_json"), [])
    fragment_text = _sanitize_field_text(
        attrs.get("python_pptx_batch_fragment", ""),
        tag_name="python_pptx_batch_fragment",
    )
    notes_text = _sanitize_field_text(
        attrs.get("pptx_batch_notes_markdown", ""),
        tag_name="pptx_batch_notes_markdown",
    )
    outputs.append(
        {
            "batch_index": batch_index + 1,
            "python_pptx_batch_fragment": fragment_text,
            "pptx_batch_notes_markdown": notes_text,
        }
    )
    return {
        "pptx_batch_index": next_index,
        "pptx_batch_iteration_ready": next_index >= batch_count,
        "pptx_batch_sections_json": attrs.get("pptx_batch_sections_json", "[]"),
        "pptx_batch_size": int(attrs.get("pptx_batch_size") or get_stage_batch_size(default=10)),
        "pptx_batch_count": batch_count,
        "pptx_batch_outputs_json": json.dumps(outputs, ensure_ascii=False),
    }


def finalize_pptx_review(_: dict[str, object], attrs: dict[str, object]) -> dict[str, object]:
    recreation_success = parse_bool(attrs.get("pptx_recreation_success"))
    human_ready = parse_bool(attrs.get("pptx_human_decision"))
    ready = human_ready if recreation_success else False

    feedback = ""
    if not ready:
        feedback_parts = [
            "Execution summary:\n" + str(attrs.get("recreated_pptx_summary_markdown", "")).strip(),
        ]
        human_feedback = str(attrs.get("pptx_human_feedback", "")).strip()
        if human_feedback:
            feedback_parts.append("Human review:\n" + human_feedback)
        feedback = "\n\n".join(part for part in feedback_parts if part.strip()).strip()

    report = "\n".join(
        [
            "# PPTX Review",
            "",
            f"- Recreation success: {recreation_success}",
            f"- Human ready: {human_ready}",
            f"- Final ready: {ready}",
            (
                "- Gate policy: on recreation failure, loop directly with execution feedback; "
                "on recreation success, human review is authoritative."
            ),
            "",
            "## Validation PDF Summary",
            str(attrs.get("recreated_pptx_summary_markdown", "")),
            "",
            "## Human Feedback",
            str(attrs.get("pptx_human_feedback", "")) or "(none)",
        ]
    )
    review_path = write_text(attrs["pptx_review_path"], report)
    return {
        "pptx_ready": ready,
        "pptx_iteration_feedback": feedback,
        "pptx_review_path": review_path,
        "outline_markdown": attrs.get("outline_markdown", ""),
        "html_content": attrs.get("html_content", ""),
        "page_budget": attrs.get("page_budget", ""),
        "recreated_script_path": attrs.get("recreated_script_path", ""),
        "recreated_pptx_path": attrs.get("recreated_pptx_path", ""),
        "recreated_validation_pdf_path": attrs.get("recreated_validation_pdf_path", ""),
        "recreated_preview_dir": attrs.get("recreated_preview_dir", ""),
        "recreated_pptx_summary_markdown": attrs.get("recreated_pptx_summary_markdown", ""),
        "recreation_notes_path": attrs.get("recreation_notes_path", ""),
        "python_pptx_script": attrs.get("python_pptx_script", ""),
        "recreation_notes_markdown": attrs.get("recreation_notes_markdown", ""),
    }


def should_terminate_pptx(_: dict[str, object], attrs: dict[str, object]) -> bool:
    return parse_bool(attrs.get("pptx_ready"))


def should_terminate_pptx_batch_loop(_: dict[str, object], attrs: dict[str, object]) -> bool:
    return parse_bool(attrs.get("pptx_batch_iteration_ready"))
__all__ = [
    "assemble_pptx_script",
    "finalize_pptx_batch_iteration",
    "finalize_pptx_review",
    "pack_pptx_review_for_human",
    "prepare_pptx_batch_iteration",
    "prepare_pptx_batch_plan",
    "run_recreation_script",
    "should_terminate_pptx",
    "should_terminate_pptx_batch_loop",
]
